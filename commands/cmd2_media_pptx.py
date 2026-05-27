import subprocess,os,datetime,json,shutil,urllib.request,base64

MP="/mnt/vnt-data/FileServer/VNT_World_AI_Division/MemPalace.md"
PROJECT_DIR="/mnt/vnt-data/FileServer/VNT_World_AI_Division/Projects/BirdHouse_Project"
WEB_GEN="/home/k/vnt-web/generated"
M4="192.168.10.94"
NC="http://192.168.10.10"
NC_ADMIN="administrator"
NC_PASS="0568116899"

def save(e):
    ts=datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
    try: open(MP,"a").write(chr(10)+"### Media ["+ts+"]"+chr(10)+e+chr(10))
    except: pass

os.makedirs(PROJECT_DIR+"/media",exist_ok=True)
os.makedirs(WEB_GEN,exist_ok=True)
os.makedirs(PROJECT_DIR+"/reports",exist_ok=True)

print("=== POWERPOINT ===")
subprocess.run(["pip","install","python-pptx","--break-system-packages","-q"],capture_output=True)
pptx_path=PROJECT_DIR+"/reports/BirdHouse_Presentation.pptx"
try:
    from pptx import Presentation
    from pptx.util import Inches,Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    prs=Presentation()
    prs.slide_width=Inches(13.33)
    prs.slide_height=Inches(7.5)
    def aslide(): return prs.slides.add_slide(prs.slide_layouts[6])
    def sbg(sl,r,g,b):
        f=sl.background.fill; f.solid(); f.fore_color.rgb=RGBColor(r,g,b)
    def atxt(sl,tx,l,t,w,h,sz=20,bd=False,cl=(255,255,255),al=PP_ALIGN.LEFT):
        tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=al
        rn=p.add_run(); rn.text=tx
        rn.font.size=Pt(sz); rn.font.bold=bd; rn.font.color.rgb=RGBColor(*cl)
    N=chr(10)
    slides=[
        {"t":"Community Bird Sanctuary","s":"VNT World AI Division  |  Ryan Khawaja  |  2026","bg":(10,30,10)},
        {"t":"Project Overview","c":"Location: Riyadh, Saudi Arabia"+N+"Scale: 100m x 68m (Football field equivalent)"+N+"Client: Ryan Khawaja"+N+"Goal: Affordable community bird sanctuary"+N+"Team: 8 VNT AI Specialists"+N+"Status: Ready for Approval","bg":(15,35,15)},
        {"t":"Site Layout Design","c":"Central Main Tower (8m radius) - focal point"+N+"4 Corner Towers (5m radius) - satellite stations"+N+"Observation Paths: Cross layout 100m x 68m"+N+"Central Pond with water bird features"+N+"Native tree planting zones (3 locations)"+N+"Main entrance with welcome signage","bg":(18,40,18)},
        {"t":"Bird Species (Saudi Arabia)","c":"Laughing Dove / Pigeon family - most common"+N+"Arabian Babbler - social native species"+N+"White-eared Bulbul - colorful resident"+N+"Little Green Bee-eater - vibrant visitor"+N+"Common Kingfisher - pond birds"+N+"Purple Sunbird - nectar feeders","bg":(15,35,25)},
        {"t":"Financial Summary","c":"Construction: SAR 180,000 - 250,000"+N+"Annual Operations: SAR 35,000/year"+N+"Funding: Municipality grants + CSR sponsors"+N+"Timeline: 6-8 months to completion"+N+"Community ROI: High social & tourism value"+N+"Analyst: Maya (VNT Finance)","bg":(20,40,30)},
        {"t":"Legal & Compliance","c":"Building permit from municipality"+N+"Environmental impact clearance"+N+"Public space authorization"+N+"Saudi environmental law compliance"+N+"Insurance and liability framework"+N+"Advisor: Amr (VNT Legal)","bg":(15,30,40)},
        {"t":"Marketing Strategy","c":"Brand: SkyHaven Community Sanctuary"+N+'Tagline: "Where Wings Find Home"'+N+"Channels: Instagram, Twitter, WhatsApp"+N+"Events: School visits, birdwatching clubs"+N+"Media: Photos, videos, 3D virtual tour"+N+"Manager: Lee (VNT Marketing)","bg":(30,20,40)},
        {"t":"Project Timeline","c":"Months 1-2: Permits & Final Design"+N+"Month 3: Site Preparation & Clearing"+N+"Months 4-5: Tower & Structure Build"+N+"Month 6: Landscaping & Water Features"+N+"Month 7: Bird Introduction Program"+N+"Month 8: Grand Opening Ceremony","bg":(40,20,20)},
        {"t":"VNT AI Division Team","c":"Alias (CEO & Lead) | Julian (Project Manager)"+N+"Nova (Civil Architect) | Ethan (Environment)"+N+"Lee (Marketing) | Amr (Legal Advisor)"+N+"Maya (Finance) | Ava (Documentation)"+N+"All work logged in MemPalace & Nextcloud","bg":(25,20,40)},
        {"t":"Next Steps","c":"1. Review and approve this proposal"+N+"2. Confirm the budget framework"+N+"3. Apply for municipal permits"+N+"4. Schedule site survey and visit"+N+"5. Sign off on DXF architectural drawings"+N+N+"Contact: kraheelw@yahoo.com | +966568116899","bg":(15,40,15)},
    ]
    for sd in slides:
        sl=aslide(); sbg(sl,*sd["bg"])
        atxt(sl,sd["t"],0.5,0.2,12.3,1.3,34,True,(80,255,80),PP_ALIGN.CENTER)
        if "s" in sd: atxt(sl,sd["s"],0.5,1.7,12.3,0.8,17,False,(160,255,160),PP_ALIGN.CENTER)
        if "c" in sd: atxt(sl,sd["c"],0.8,1.8,11.7,5.1,17,False,(210,255,210))
        atxt(sl,"VNT World AI Division  |  Confidential  |  Prepared for Ryan Khawaja",0.5,7.0,12.3,0.4,9,False,(60,120,60),PP_ALIGN.CENTER)
    prs.save(pptx_path)
    shutil.copy(pptx_path,WEB_GEN+"/BirdHouse_Presentation.pptx")
    print("PowerPoint OK:", os.path.getsize(pptx_path)//1024,"KB")
    save("PowerPoint created: "+pptx_path)
except Exception as e: print("PowerPoint error:",str(e))

print(chr(10)+"=== REQUESTING IMAGES ===")
prompts=[
    ("birdhouse_aerial_view","aerial view beautiful community bird sanctuary football field size multiple tall wooden towers birds flying lush green trees landscape sunset"),
    ("birdhouse_main_tower","majestic tall wooden community bird tower sanctuary birds perching natural materials beautiful architecture golden hour"),
    ("birdhouse_family_visit","happy families children parents visiting bird sanctuary observation deck watching colorful birds feeding joyful community park"),
    ("birdhouse_pond_birds","peaceful bird sanctuary pond water feature ducks pigeons doves birds drinking splashing surrounded by trees beautiful nature"),
]
for name,prompt in prompts:
    try:
        body=json.dumps({"prompt":prompt,"width":1024,"height":768,"steps":25}).encode()
        req=urllib.request.Request("http://"+M4+":3333/generate",data=body,headers={"Content-Type":"application/json"},method="POST")
        with urllib.request.urlopen(req,timeout=90) as r: d=json.loads(r.read())
        path=d.get("path","") or ""
        if path and os.path.exists(path):
            shutil.copy(path,WEB_GEN+"/"+name+".png")
            print("Image saved:",name)
        else: print("Image queued:",name,d.get("error",""))
    except Exception as e: print("Image",name,":",str(e))

print(chr(10)+"=== REQUESTING VIDEO ===")
try:
    body=json.dumps({"prompt":"cinematic aerial tour bird sanctuary towers birds flying golden hour beautiful landscape","frames":24,"fps":8,"steps":20}).encode()
    req=urllib.request.Request("http://"+M4+":3333/generate-video",data=body,headers={"Content-Type":"application/json"},method="POST")
    with urllib.request.urlopen(req,timeout=120) as r: d=json.loads(r.read())
    print("Video:",d.get("url") or d.get("path") or str(d))
except Exception as e: print("Video:",str(e))

print(chr(10)+"=== REQUESTING 3D ===")
try:
    body=json.dumps({"description":"bird sanctuary wooden tower multiple perch levels birds natural","format":"obj","render_preview":True}).encode()
    req=urllib.request.Request("http://"+M4+":3333/generate-3d",data=body,headers={"Content-Type":"application/json"},method="POST")
    with urllib.request.urlopen(req,timeout=120) as r: d=json.loads(r.read())
    print("3D:",d.get("url") or d.get("path") or str(d))
except Exception as e: print("3D:",str(e))

print(chr(10)+"=== UPLOADING TO NEXTCLOUD ===")
auth_hdr=base64.b64encode((NC_ADMIN+":"+NC_PASS).encode()).decode()
def nc_up(local,remote_name):
    if not os.path.exists(local): return
    try:
        with open(local,"rb") as f: data=f.read()
        url=NC+"/remote.php/dav/files/"+NC_ADMIN+"/VNT_Projects/BirdHouse/"+remote_name
        req=urllib.request.Request(url,data=data,headers={"Authorization":"Basic "+auth_hdr,"Content-Type":"application/octet-stream"},method="PUT")
        with urllib.request.urlopen(req,timeout=30) as r: print("NC uploaded:",remote_name,r.status)
    except Exception as e: print("NC",remote_name,":",str(e))

nc_up(pptx_path,"BirdHouse_Presentation.pptx")
nc_up(WEB_GEN+"/bird_house_proposal.html","bird_house_proposal.html")
nc_up(WEB_GEN+"/birdhouse_site_plan.dxf","birdhouse_site_plan.dxf")
if os.path.exists(PROJECT_DIR+"/documents"):
    for doc in os.listdir(PROJECT_DIR+"/documents"):
        nc_up(PROJECT_DIR+"/documents/"+doc, doc)

save("Media requested from M4. PowerPoint done. Nextcloud upload attempted.")
print("CMD2 COMPLETE")
